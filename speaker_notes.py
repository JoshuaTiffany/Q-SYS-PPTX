import socket
from pptx import Presentation
import json
import re

presentation_path = r"C:\Users\Joshua.Tiffany\Downloads\Intern Workshop - Financial Literacy 06.21.23.pptx"

def extract_speaker_notes(presentation_path):
    presentation = Presentation(presentation_path)
    notes = []

    for slide in presentation.slides:
        slide_number_value = slide_number(presentation.slides, slide)
        speaker_notes = []
        info_data = []  # Variable to store info data
        has_video= False
        is_auto = False

        for shape in slide.notes_slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "[video]" in run.text:
                            has_video = True
                        elif "[auto]" in run.text:
                            is_auto = True
                        elif re.search(r"[Ii]nfo\[\s*(.*?)\s*\]", run.text):
                            matches = re.findall(r"[Ii]nfo\[\s*(.*?)\s*\]", run.text)
                            if matches:
                                info_data.extend(matches)
                        else:
                            speaker_notes.append(run.text)

        slide_data = {
            "slide_number": slide_number_value,
            "speaker_notes": speaker_notes,
            "info_data": info_data  # Include info data in slide data
        }

        if has_video:
            slide_data["video"] = True

        if is_auto:
            slide_data["auto"] = True

        notes.append(slide_data)

    return notes


def slide_number(slides, slide):
    for idx, s in enumerate(slides):
        if s == slide:
            return idx + 1


# Set up TCP server
server_socket = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
server_address = ('localhost', 12345)
server_socket.bind(server_address)

print('Server listening on {}:{}'.format(*server_address))

# Listen for client connection
server_socket.listen(1)
print('Waiting for client connection...')

# Accept client connection
client_socket, client_address = server_socket.accept()
print('Client connected:', client_address)

# Extract speaker notes from the PowerPoint presentation
notes = extract_speaker_notes(presentation_path)

# Send slide data to the client
for slide in notes:
    slide_number = slide["slide_number"]
    speaker_notes = slide["speaker_notes"]
    has_video = slide.get("video", False)  # Check if "video" key exists and set default value to False
    is_auto = slide.get("auto", False)  # Check if "video" key exists and set default value to False
    info_data = slide["info_data"]

    video_status = "True" if has_video else "False"
    auto_status = "True" if is_auto else "False"
    data = f"Slide {slide_number}: {speaker_notes}\r\n "
    
    
    info = f" Info Data: {', '.join(info_data)}\r\n "
    
  
    video = f" Video Data: {video_status}\r\n "
    auto = f" Auto Data: {auto_status}\r\n "
    
    client_socket.sendall(data.encode())
    client_socket.sendall(info.encode())
    client_socket.sendall(video.encode())
    client_socket.sendall(auto.encode())

# Close the connection and socket 
client_socket.close()
server_socket.close()


